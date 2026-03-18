import streamlit as st
import pandas as pd
import numpy as np
import matplotlib.pyplot as plt
from scipy import stats
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from io import BytesIO

# --- Configuration & Theme ---
st.set_page_config(layout="wide", page_title="Performance Headroom Calculator")

VIKING_BLUE = RGBColor(0, 107, 230)
TEXT_GREY = RGBColor(89, 89, 89)
LIGHT_GREY = RGBColor(166, 166, 166)

# --- PowerPoint Generation Function ---
def create_viking_slide(client, vertical, cpa_goal, daily_spend, total_budget, days, chart_buffer):
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6]) # Blank Layout

    # 1. Top Left: Source [cite: 1]
    source_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(4), Inches(0.5))
    sp = source_box.text_frame.paragraphs[0]
    sp.text = "Source: Performance Headroom Tool Analysis"
    sp.font.size = Pt(10)
    sp.font.color.rgb = LIGHT_GREY

    # 2. Main Headline [cite: 4]
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(1))
    title_box.text_frame.word_wrap = True
    tp = title_box.text_frame.paragraphs[0]
    tp.text = f"{client if client else 'Client'} can support a {days}-day budget of ${total_budget:,.0f} at a ${cpa_goal:,.0f} CPA."
    tp.font.size = Pt(22)
    tp.font.bold = True

    # 3. Headroom Label & Blue Arrow [cite: 3]
    arrow = slide.shapes.add_shape(MSO_SHAPE.UP_ARROW, Inches(0.45), Inches(1.35), Inches(0.3), Inches(0.4))
    arrow.fill.solid()
    arrow.fill.fore_color.rgb = VIKING_BLUE
    arrow.line.fill.background()

    headroom_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.3), Inches(2), Inches(0.5))
    hp = headroom_box.text_frame.paragraphs[0]
    hp.text = "Headroom"
    hp.font.size = Pt(18)
    hp.font.color.rgb = VIKING_BLUE
    
    # 4. Main Chart Placement [cite: 8]
    chart_buffer.seek(0)
    slide.shapes.add_picture(chart_buffer, Inches(0.5), Inches(2.2), width=Inches(5.8))

    # 5. Right Side Analysis Info [cite: 7]
    label_box = slide.shapes.add_textbox(Inches(6.5), Inches(2.2), Inches(3), Inches(0.5))
    lp = label_box.text_frame.paragraphs[0]
    lp.text = "Headroom Regression Analysis (Linear)"
    lp.font.size = Pt(12)
    lp.font.bold = True
    lp.font.color.rgb = TEXT_GREY

    meta_box = slide.shapes.add_textbox(Inches(6.5), Inches(2.6), Inches(3), Inches(0.5))
    mp = meta_box.text_frame.paragraphs[0]
    mp.text = vertical.upper() if vertical else "CAMPAIGN_STRATEGY"
    mp.font.size = Pt(10)
    mp.font.color.rgb = VIKING_BLUE

    # 6. Bottom Takeaway [cite: 6]
    takeaway_box = slide.shapes.add_textbox(Inches(0.5), Inches(6.8), Inches(9), Inches(0.5))
    tkp = takeaway_box.text_frame.paragraphs[0]
    tkp.text = f"Takeaway: Recommended daily spend of ${daily_spend:,.2f} to maintain efficiency goal."
    tkp.font.size = Pt(14)
    tkp.font.bold = True

    binary_output = BytesIO()
    prs.save(binary_output)
    return binary_output.getvalue()

# --- Streamlit UI ---
st.title('📈 Performance Headroom Calculator')

st.header("Step 1: Campaign Information")
col_meta1, col_meta2 = st.columns(2)
with col_meta1:
    client_name = st.text_input("Client Name", placeholder="e.g., Viking River Cruises")
with col_meta2:
    vertical_name = st.text_input("Vertical", placeholder="e.g., River Cruises")

st.write("---")
st.header("Step 2: Upload Your Data")
uploaded_file = st.file_uploader("Choose a CSV or Excel file", type=['csv', 'xlsx'])

if uploaded_file:
    try:
        df = pd.read_csv(uploaded_file) if uploaded_file.name.endswith('.csv') else pd.read_excel(uploaded_file)
        
        required_cols = ['Date', 'Revenue (USD)', 'Total Conversions']
        if not all(col in df.columns for col in required_cols):
            st.error(f"Missing columns: {required_cols}")
        else:
            # Data Cleaning & Analysis
            df['Date'] = pd.to_datetime(df['Date'])
            df_daily = df.groupby('Date').agg({'Revenue (USD)': 'sum', 'Total Conversions': 'sum'}).reset_index()
            df_daily['eCPA'] = df_daily['Revenue (USD)'] / df_daily['Total Conversions']
            df_daily.replace([np.inf, -np.inf], np.nan, inplace=True)

            # Outlier Detection
            q1, q3 = df_daily['eCPA'].quantile(0.25), df_daily['eCPA'].quantile(0.75)
            iqr = q3 - q1
            df_daily['is_outlier'] = (df_daily['eCPA'] < (q1 - 1.5*iqr)) | (df_daily['eCPA'] > (q3 + 1.5*iqr))
            df_cleaned = df_daily[~df_daily['is_outlier']].dropna()

            # Regression
            slope, intercept, r_value, _, _ = stats.linregress(df_cleaned['Revenue (USD)'], df_cleaned['eCPA'])
            st.session_state.model = {'slope': slope, 'intercept': intercept}

            # Plotting
            fig, ax = plt.subplots(figsize=(10, 5))
            ax.scatter(df_daily['Revenue (USD)'], df_daily['eCPA'], c=df_daily['is_outlier'], cmap='coolwarm', alpha=0.6)
            x_vals = np.array(ax.get_xlim())
            ax.plot(x_vals, intercept + slope * x_vals, '--', color='red', label='Trendline')
            ax.set_title(f"Spend vs CPA: {client_name}")
            st.pyplot(fig)

            # Store chart for PPTX
            chart_buf = BytesIO()
            fig.savefig(chart_buf, format='png', dpi=300)
            st.session_state.chart_buf = chart_buf

    except Exception as e:
        st.error(f"Error: {e}")

st.write("---")
st.subheader("Step 3: Calculate & Export")
if 'model' in st.session_state:
    c1, c2 = st.columns(2)
    goal_cpa = c1.number_input("Goal CPA ($)", min_value=0.0, value=1000.0)
    days = c2.number_input("Forecast Days", min_value=1, value=30)

    if st.button("Generate Recommendation"):
        m = st.session_state.model
        daily_spend = (goal_cpa - m['intercept']) / m['slope']
        total_budget = daily_spend * days
        
        st.metric("Recommended Daily Spend", f"${daily_spend:,.2f}")
        st.metric("Total Period Budget", f"${total_budget:,.2f}")

        pptx_file = create_viking_slide(client_name, vertical_name, goal_cpa, daily_spend, total_budget, days, st.session_state.chart_buf)
        st.download_button("📥 Download Viking Forecast Slide", pptx_file, f"{client_name}_Headroom.pptx")
